import io
import os
import re
import asyncio
from pathlib import Path
from typing import Optional
import structlog

logger = structlog.get_logger()


class FileExtractor:
    """Extracts raw text from uploaded files using appropriate parsers."""

    SUPPORTED_IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".webp", ".bmp", ".tiff"}
    SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".ppt"} | SUPPORTED_IMAGE_EXTENSIONS

    async def extract(self, file_bytes: bytes, filename: str) -> str:
        """Main entry point. Dispatches to correct extractor by file extension."""
        ext = Path(filename).suffix.lower()

        if ext == ".pdf":
            return await self._extract_pdf(file_bytes)
        elif ext == ".docx":
            return await self._extract_docx(file_bytes)
        elif ext in {".pptx", ".ppt"}:
            return await self._extract_pptx(file_bytes)
        elif ext in self.SUPPORTED_IMAGE_EXTENSIONS:
            return await self._extract_image_ocr(file_bytes)
        else:
            raise ValueError(f"Unsupported file type: {ext}")

    async def _extract_pdf(self, file_bytes: bytes) -> str:
        """Extract text from PDF using PyMuPDF. Falls back to OCR for scanned pages."""
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(None, self._extract_pdf_sync, file_bytes)

    def _extract_pdf_sync(self, file_bytes: bytes) -> str:
        import fitz  # PyMuPDF

        text_pages = []
        doc = fitz.open(stream=file_bytes, filetype="pdf")

        for page_num, page in enumerate(doc):
            text = page.get_text("text").strip()

            if len(text) < 50:
                # Likely scanned — attempt OCR on page image
                pix = page.get_pixmap(dpi=200)
                img_bytes = pix.tobytes("png")
                text = self._ocr_image_bytes(img_bytes)
                logger.info("PDF page OCR fallback", page=page_num + 1)

            if text:
                text_pages.append(f"[Page {page_num + 1}]\n{text}")

        doc.close()
        full_text = "\n\n".join(text_pages)
        return self._clean_text(full_text)

    async def _extract_docx(self, file_bytes: bytes) -> str:
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(None, self._extract_docx_sync, file_bytes)

    def _extract_docx_sync(self, file_bytes: bytes) -> str:
        from docx import Document

        doc = Document(io.BytesIO(file_bytes))
        paragraphs = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                if para.style.name.startswith("Heading"):
                    paragraphs.append(f"\n## {text}")
                else:
                    paragraphs.append(text)

        # Also extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    paragraphs.append(row_text)

        return self._clean_text("\n".join(paragraphs))

    async def _extract_pptx(self, file_bytes: bytes) -> str:
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(None, self._extract_pptx_sync, file_bytes)

    def _extract_pptx_sync(self, file_bytes: bytes) -> str:
        from pptx import Presentation
        from pptx.util import Pt

        prs = Presentation(io.BytesIO(file_bytes))
        slides_text = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = [f"[Slide {slide_num}]"]

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_content.append(text)

            if len(slide_content) > 1:
                slides_text.append("\n".join(slide_content))

        return self._clean_text("\n\n".join(slides_text))

    async def _extract_image_ocr(self, file_bytes: bytes) -> str:
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(None, self._extract_image_ocr_sync, file_bytes)

    def _extract_image_ocr_sync(self, file_bytes: bytes) -> str:
        from PIL import Image, ImageEnhance, ImageFilter
        import pytesseract

        img = Image.open(io.BytesIO(file_bytes))

        # Preprocessing for better OCR accuracy
        img = img.convert("L")  # Grayscale
        img = ImageEnhance.Contrast(img).enhance(2.0)
        img = img.filter(ImageFilter.SHARPEN)

        # Resize if too small
        width, height = img.size
        if width < 1000:
            scale = 1000 / width
            img = img.resize((int(width * scale), int(height * scale)), Image.LANCZOS)

        config = "--oem 3 --psm 6 -l eng"
        text = pytesseract.image_to_string(img, config=config)
        return self._clean_text(text)

    def _ocr_image_bytes(self, image_bytes: bytes) -> str:
        """OCR from raw bytes (used for scanned PDF pages)."""
        from PIL import Image
        import pytesseract

        img = Image.open(io.BytesIO(image_bytes)).convert("L")
        config = "--oem 3 --psm 6"
        return pytesseract.image_to_string(img, config=config)

    def _clean_text(self, text: str) -> str:
        """Normalise whitespace, remove garbled characters, deduplicate blank lines."""
        text = re.sub(r'\x00', '', text)
        text = re.sub(r'[ \t]+', ' ', text)
        text = re.sub(r'\n{3,}', '\n\n', text)
        text = re.sub(r'[^\x09\x0A\x0D\x20-\x7E\u00A0-\uFFFF]', '', text)
        lines = [line.strip() for line in text.splitlines()]
        lines = [l for l in lines if l]
        return "\n".join(lines).strip()


file_extractor = FileExtractor()
